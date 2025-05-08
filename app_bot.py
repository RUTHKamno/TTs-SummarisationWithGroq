import streamlit as st
import os
import requests
from transformers import pipeline
import json
from dotenv import load_dotenv
from pathlib import Path
import base64
import io
import numpy as np
import time

# Import libraries for document generation
try:
    from pptx import Presentation
    from docx import Document
    from docx.shared import Inches
except ImportError:
    st.error("Please install python-pptx and python-docx: pip install python-pptx python-docx")


# Note: IPython.display.Audio is for Jupyter notebooks, we'll use st.audio for Streamlit

# --- Load environment variables ---
# This will load variables from a .env file if it exists
load_dotenv()

# --- Get API Keys ---
# Prioritize Streamlit secrets if available and keys are present,
# otherwise use environment variables, and finally the sidebar input.
groq_api_key = None
hf_api_key = None

# Check if st.secrets is available before accessing it
if hasattr(st, 'secrets') and st.secrets:
    # Check if keys exist in secrets before using them
    groq_api_key = st.secrets.get("GROQ_API_KEY")
    hf_api_key = st.secrets.get("HF_API_KEY")

# Fallback to environment variables if not found in secrets or secrets not available
if not groq_api_key:
    groq_api_key = os.getenv("GROQ_API_KEY")
if not hf_api_key:
    hf_api_key = os.getenv("HF_API_key") # Corrected variable name


# --- Define the Summarization Service Class ---
# (Copied and slightly adapted from your notebook)
class SummarizationService:
    def __init__(self, api_choice="huggingface", groq_api_key=None, hf_model_name="facebook/bart-large-cnn"):
        """
        Initialize summarization service with either Groq API or Hugging Face.

        Args:
            api_choice: "groq" or "huggingface"
            groq_api_key: Your Groq API key (required if using Groq)
            hf_model_name: Name of the Hugging Face model to use
        """
        self.api_choice = api_choice

        if api_choice == "groq":
            if not groq_api_key:
                # In Streamlit, we'll handle this check in the main app logic/sidebar
                pass
            self.groq_api_key = groq_api_key
            self.groq_api_url = "https://api.groq.com/openai/v1/chat/completions"

        elif api_choice == "huggingface":
            # Initialize the Hugging Face summarization pipeline
            # Streamlit caching can help here to avoid reloading the model
            @st.cache_resource
            def load_hf_summarizer(model_name):
                 return pipeline("summarization", model=model_name)

            self.summarizer = load_hf_summarizer(hf_model_name)
        else:
            raise ValueError("api_choice must be either 'groq' or 'huggingface'")

    def summarize_with_groq(self, text, max_length=150):
        """Summarize text using Groq API"""
        if not self.groq_api_key:
             return "Groq API key is not set. Please provide it in the sidebar or environment variables."

        headers = {
            "Authorization": f"Bearer {self.groq_api_key}",
            "Content-Type": "application/json"
        }

        data = {
            "model": "llama3-70b-8192", # You can change to other models Groq supports
            "messages": [
                {"role": "system", "content": "You are a helpful assistant that summarizes text."},
                {"role": "user", "content": f"Please provide a concise summary of the following text in around {max_length} words: \n\n{text}"}
            ],
            "max_tokens": 1024
        }

        try:
            response = requests.post(self.groq_api_url, headers=headers, json=data)
            response.raise_for_status() # Raise an exception for bad status codes
            response_json = response.json()

            if 'choices' in response_json and len(response_json['choices']) > 0:
                return response_json['choices'][0]['message']['content'].strip()
            else:
                return f"Groq API did not return a summary. Response: {response_json}"
        except requests.exceptions.RequestException as e:
            return f"Error communicating with Groq API: {e}"
        except Exception as e:
            return f"An unexpected error occurred during Groq summarization: {e}"


    def summarize_with_huggingface(self, text, max_length=150, min_length=50):
        """Summarize text using Hugging Face pipeline"""
        # Calculate max and min length in tokens (approximately)
        # Adjusting token calculation for better control
        words_per_token = 0.75 # Approximation
        max_tokens = int(max_length / words_per_token)
        min_tokens = int(min_length / words_per_token)

        # Ensure min_tokens is not greater than max_tokens and both are reasonable
        min_tokens = min(min_tokens, max_tokens - 10) # Ensure some difference
        min_tokens = max(min_tokens, 10) # Minimum reasonable length

        try:
            summary = self.summarizer(text,
                                      max_length=max_tokens,
                                      min_length=min_tokens,
                                      do_sample=False)

            return summary[0]['summary_text']
        except Exception as e:
            return f"An error occurred during Hugging Face summarization: {e}"


    def summarize(self, text, max_length=150):
        """Summarize text using the selected API"""
        if not text:
            return "Please enter text to summarize."
        if self.api_choice == "groq":
            # API key check is now inside summarize_with_groq
            return self.summarize_with_groq(text, max_length)
        else:
            # Hugging Face does not require a separate API key for the default model
            return self.summarize_with_huggingface(text, max_length)

# --- Define the Text-to-Speech Service Class ---
# (Copied and slightly adapted from your notebook)
class TextToSpeechService:
    def __init__(self, groq_api_key=None):
        """
        Initialize the Text-to-Speech service using Groq for both text processing and speech synthesis.

        Args:
            groq_api_key: Your Groq API key
        """
        self.groq_api_key = groq_api_key

        if not self.groq_api_key:
            # In Streamlit, we'll handle this check in the main app logic
            pass

        # Default Groq model for text processing
        self.llm_model = "llama3-70b-8192"

        # Default Groq TTS model
        self.tts_model = "playai-tts"

        # Default voice
        self.voice = "Cillian-PlayAI"

        # Available PlayAI voices
        self.playai_voices = [\
            "Aaliyah-PlayAI", "Adelaide-PlayAI", "Angelo-PlayAI", "Arista-PlayAI", \
            "Atlas-PlayAI", "Basil-PlayAI", "Briggs-PlayAI", "Calum-PlayAI", \
            "Celeste-PlayAI", "Cheyenne-PlayAI", "Chip-PlayAI", "Cillian-PlayAI", \
            "Deedee-PlayAI", "Eleanor-PlayAI", "Fritz-PlayAI", "Gail-PlayAI", \
            "Indigo-PlayAI", "Jennifer-PlayAI", "Judy-PlayAI", "Mamaw-PlayAI", \
            "Mason-PlayAI", "Mikail-PlayAI", "Mitch-PlayAI", "Nia-PlayAI", \
            "Quinn-PlayAI", "Ruby-PlayAI", "Thunder-PlayAI"\
        ]

    def set_tts_model(self, model_id):
        """
        Set a specific Groq TTS model.

        Args:
            model_id: ID of the Groq TTS model
        """
        self.tts_model = model_id

    def set_voice(self, voice):
        """
        Set a specific voice for TTS.

        Args:
            voice: Name of the voice to use

        Raises:
            ValueError: If using PlayAI model and the voice is not in the list of available voices
        """
        if self.tts_model == "playai-tts" and voice not in self.playai_voices:
            available_voices = "\\n".join(self.playai_voices)
            raise ValueError(f"The specified voice '{voice}' is not available for the PlayAI model. Available voices are:\\n{available_voices}")

        self.voice = voice

    def process_text(self, text):
        """
        Process text with Groq to enhance it for TTS (adjust tone, emphasis, etc.)

        Args:
            text: Text to be processed

        Returns:
            Enhanced text for TTS
        """
        if not self.groq_api_key:
            return text # Cannot enhance without API key

        headers = {
            "Authorization": f"Bearer {self.groq_api_key}",
            "Content-Type": "application/json"
        }

        # Prompt to enhance text for TTS
        prompt = f"""
        You are a text-to-speech enhancement assistant.
        Your task is to enhance the following text for better speech synthesis quality.
        Add appropriate punctuation for better pauses, adjust emphasis where needed,
        and format the text in a way that would sound natural when read aloud.

        Text to enhance: "{text}"

        Provide ONLY the enhanced text without any explanations or additional comments.
        """

        payload = {
            "model": self.llm_model,
            "messages": [{"role": "user", "content": prompt}],
            "temperature": 0.2,
            "max_tokens": 1024
        }

        try:
            response = requests.post(
                "https://api.groq.com/openai/v1/chat/completions",
                headers=headers,
                json=payload
            )

            if response.status_code == 200:
                enhanced_text = response.json()["choices"][0]["message"]["content"].strip()
                return enhanced_text
            else:
                st.warning(f"Error processing text with Groq for enhancement: {response.status_code}, {response.text}")
                return text # Return original text if enhancement fails
        except requests.exceptions.RequestException as e:
            st.warning(f"Error communicating with Groq API for text enhancement: {e}")
            return text
        except Exception as e:
            st.warning(f"An unexpected error occurred during Groq text enhancement: {e}")
            return text


    def generate_speech(self, text, enhance_text=True):
        """
        Generate speech from text using Groq's TTS capabilities.

        Args:
            text: Text to convert to speech
            enhance_text: Whether to enhance the text with Groq first

        Returns:
            Audio data bytes or None if generation fails
        """
        if not self.groq_api_key:
            st.error("Groq API key is not set. Cannot generate speech.")
            return None

        if not text:
            st.warning("Please enter text to convert to speech.")
            return None

        if enhance_text:
            text = self.process_text(text)
            # st.info(f"Enhanced text for TTS: {text}") # Optional: show enhanced text

        # Configure Groq API call
        headers = {
            "Authorization": f"Bearer {self.groq_api_key}",
            "Content-Type": "application/json"
        }

        # Prepare payload for Groq API
        payload = {
            "model": self.tts_model,
            "input": text,
            "voice": self.voice,
            "response_format": "mp3"
        }

        # Make the API call to Groq TTS API endpoint
        api_url = "https://api.groq.com/openai/v1/audio/speech"

        # Implement exponential backoff for API calls
        max_retries = 5
        retry_delay = 1

        for attempt in range(max_retries):
            try:
                response = requests.post(api_url, headers=headers, json=payload)

                if response.status_code == 200:
                    # Success
                    return response.content
                elif response.status_code == 503 or response.status_code == 429:
                    # Model is loading or rate limit
                    st.warning(f"Service busy or rate limited, retrying in {retry_delay} seconds... (Attempt {attempt + 1}/{max_retries})")
                    time.sleep(retry_delay)
                    retry_delay *= 2  # Exponential backoff
                else:
                    st.error(f"Error calling Groq TTS API: {response.status_code}, {response.text}")
                    return None
            except requests.exceptions.RequestException as e:
                st.error(f"Error during API call: {str(e)}")
                time.sleep(retry_delay)
                retry_delay *= 2

        st.error("Failed to generate speech after multiple attempts")
        return None

# --- Chatbot Functionality ---

def get_groq_response(prompt, api_key, model="llama3-8b-8192"):
    """Gets a response from the Groq API."""
    if not api_key:
        return "Groq API key is not set. Please provide it in the sidebar or environment variables."

    headers = {
        "Authorization": f"Bearer {api_key}",
        "Content-Type": "application/json"
    }

    data = {
        "model": model,
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": 1024,
        "temperature": 0.7,
    }

    try:
        response = requests.post("https://api.groq.com/openai/v1/chat/completions", headers=headers, json=data)
        response.raise_for_status()
        response_json = response.json()
        if 'choices' in response_json and len(response_json['choices']) > 0:
            return response_json['choices'][0]['message']['content'].strip()
        else:
            return f"Groq API did not return a response. Response: {response_json}"
    except requests.exceptions.RequestException as e:
        return f"Error communicating with Groq API: {e}"
    except Exception as e:
        return f"An unexpected error occurred during chatbot interaction: {e}"

# --- Document Export Functions ---

def create_ppt_from_chat(messages):
    """Creates a PowerPoint presentation from chat messages."""
    prs = Presentation()
    # Use a blank slide layout (layout index 6 is typically blank)
    blank_slide_layout = prs.slide_layouts[6]

    for message in messages:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(1) # Placeholder values

        # Add a text box for the message
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), prs.slide_width - Inches(1), prs.slide_height - Inches(1))
        tf = txBox.text_frame

        p = tf.add_paragraph()
        p.text = f"{message['role'].capitalize()}: {message['content']}"

        # Optionally, style based on role
        if message['role'] == 'user':
            p.font.bold = True
        else:
            p.font.italic = True

    # Save the presentation to a bytes buffer
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def create_word_from_text(text):
    """Creates a Word document from text."""
    document = Document()
    document.add_paragraph(text)

    # Save the document to a bytes buffer
    buf = io.BytesIO()
    document.save(buf)
    buf.seek(0)
    return buf


# --- Streamlit App Layout ---
st.set_page_config(page_title="Text Utility and Chatbot App", layout="centered")

st.title("Text Utility and Chatbot")

st.markdown("""
This app allows you to summarize text, convert text to speech, or chat with a Groq-powered assistant.
You can also export your chat discussions.
""")

# --- Sidebar for API Key Input and Settings ---
with st.sidebar:
    st.header("Settings")

    # API Keys input
    st.subheader("API Keys")
    st.warning("It's recommended to use Streamlit Secrets or environment variables for API keys.")
    # Use the fetched groq_api_key as the initial value for the text input
    groq_api_key_input = st.text_input("Groq API Key", type="password", value=groq_api_key)
    # Hugging Face API Key is generally not needed for the default summarization model,
    # but you could add it here if you used a different model requiring auth.
    # hf_api_key_input = st.text_input("Hugging Face API Key (Optional)", type="password", value=hf_api_key)

    # Update the groq_api_key variable if user inputs it in the sidebar
    if groq_api_key_input:
        groq_api_key = groq_api_key_input
    # if hf_api_key_input:
    #     hf_api_key = hf_api_key_input


    st.subheader("Choose Action")
    action = st.radio("Select an action:", ("Summarize Text", "Text-to-Speech", "Chatbot"))

    # Action-specific settings
    if action == "Summarize Text":
        st.subheader("Summarization Settings")
        summarizer_choice = st.radio("Choose Summarizer:", ("Hugging Face (Default)", "Groq"))
        max_summary_length = st.slider("Maximum Summary Length (words)", 50, 500, 150)
        # You could add an option for HF model name here if needed
        # hf_model = st.text_input("Hugging Face Model Name", "facebook/bart-large-cnn")

    elif action == "Text-to-Speech":
        st.subheader("Text-to-Speech Settings")
        # Initialize TTS service to get available voices
        # Only attempt to initialize if a Groq API key is available
        if groq_api_key:
            try:
                tts_service_for_voices = TextToSpeechService(groq_api_key=groq_api_key)
                available_voices = tts_service_for_voices.playai_voices
                selected_voice_tts = st.selectbox("Select Voice:", available_voices, index=available_voices.index("Cillian-PlayAI") if "Cillian-PlayAI" in available_voices else 0, key="tts_voice_select")
                enhance_text_for_tts = st.checkbox("Enhance text for better speech (uses Groq)", value=True, key="tts_enhance_checkbox")
            except ValueError as e:
                st.warning(f"Could not initialize TTS service to list voices: {e}. Please ensure Groq API key is set correctly.")
                selected_voice_tts = None
                enhance_text_for_tts = st.checkbox("Enhance text for better speech (uses Groq)", value=True, disabled=True, key="tts_enhance_checkbox_disabled") # Disable if no key
                available_voices = []
        else:
             st.warning("Enter your Groq API key above to configure Text-to-Speech settings.")
             selected_voice_tts = None
             enhance_text_for_tts = st.checkbox("Enhance text for better speech (uses Groq)", value=True, disabled=True, key="tts_enhance_checkbox_disabled_2") # Disable if no key
             available_voices = []

    elif action == "Chatbot":
        st.subheader("Chatbot Settings")
        # Chatbot specific settings can go here if needed (e.g., model choice)
        chatbot_model = st.selectbox("Groq Model:", ("llama3-8b-8192", "llama3-70b-8192"), key="chatbot_model_select")
        max_summary_length_chat = st.slider("Max Summary Length for Audio (words)", 30, 200, 80, key="chat_summary_length")
        # Voice selection for chatbot audio response
        if groq_api_key:
            try:
                tts_service_for_voices_chat = TextToSpeechService(groq_api_key=groq_api_key)
                available_voices_chat = tts_service_for_voices_chat.playai_voices
                selected_voice_chat = st.selectbox("Select Voice for Chat Audio:", available_voices_chat, index=available_voices_chat.index("Cillian-PlayAI") if "Cillian-PlayAI" in available_voices_chat else 0, key="chat_voice_select")
                enhance_text_for_tts_chat = st.checkbox("Enhance chat response text for speech (uses Groq)", value=True, key="chat_enhance_checkbox")
            except ValueError as e:
                 st.warning(f"Could not initialize TTS service to list voices for chat audio: {e}. Please ensure Groq API key is set correctly.")
                 selected_voice_chat = None
                 enhance_text_for_tts_chat = st.checkbox("Enhance chat response text for speech (uses Groq)", value=True, disabled=True, key="chat_enhance_checkbox_disabled") # Disable if no key
                 available_voices_chat = []
        else:
             st.warning("Enter your Groq API key above to configure Chatbot audio settings.")
             selected_voice_chat = None
             enhance_text_for_tts_chat = st.checkbox("Enhance chat response text for speech (uses Groq)", value=True, disabled=True, key="chat_enhance_checkbox_disabled_2") # Disable if no key
             available_voices_chat = []

    # Add a button to export the entire discussion to PowerPoint in the sidebar
    # This button is available regardless of the main action, but only if there are messages
    if "messages" in st.session_state and st.session_state.messages:
        st.subheader("Export Chat")
        ppt_file_buffer = create_ppt_from_chat(st.session_state.messages)
        st.download_button(
            label="📊 Export Full Discussion to PowerPoint",
            data=ppt_file_buffer,
            file_name="chatbot_discussion.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="export_ppt_button_sidebar" # Unique key for the download button in sidebar
        )


# --- Main Content Area ---
st.header("Input")

if action in ["Summarize Text", "Text-to-Speech"]:
    input_text = st.text_area("Enter the text here:", height=200)
    result = None
    audio_data = None

    if st.button(action):
        if not input_text:
            st.warning("Please enter some text.")
        else:
            with st.spinner(f"Performing {action.lower()}..."):
                if action == "Summarize Text":
                    api_choice = "groq" if summarizer_choice == "Groq" else "huggingface"
                    try:
                        summarization_service = SummarizationService(
                            api_choice=api_choice,
                            groq_api_key=groq_api_key
                            # hf_model_name=hf_model # Uncomment if you added model input
                        )
                        result = summarization_service.summarize(input_text, max_length=max_summary_length)
                    except ValueError as e:
                        st.error(f"Configuration Error: {e}")
                    except Exception as e:
                        st.error(f"An unexpected error occurred: {e}")

                elif action == "Text-to-Speech":
                    if not groq_api_key:
                         st.error("Groq API key is not set. Cannot generate speech.")
                    else:
                        try:
                            tts_service = TextToSpeechService(groq_api_key=groq_api_key)
                            if selected_voice_tts: # Only set voice if voices were loaded successfully
                                tts_service.set_voice(selected_voice_tts)
                            audio_data = tts_service.generate_speech(input_text, enhance_text=enhance_text_for_tts)
                        except ValueError as e:
                            st.error(f"Configuration Error: {e}")
                        except Exception as e:
                            st.error(f"An unexpected error occurred: {e}")

    # --- Display Results for Summarize/TTS ---
    st.header("Result")

    if result:
        st.subheader("Summary:")
        st.write(result)

    if audio_data:
        st.subheader("Audio Output:")
        # Use st.audio to play the generated audio
        st.audio(audio_data, format='audio/mp3')


elif action == "Chatbot":
    st.header("Chat with Groq")

    # Initialize chat history
    if "messages" not in st.session_state:
        st.session_state.messages = []

    # Display chat messages from history
    for i, message in enumerate(st.session_state.messages):
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
            # Add buttons for AI responses
            if message["role"] == "assistant":
                # Use a unique key for each button based on the message index
                button_key_vocalize = f"vocalize_btn_{i}"
                button_key_word = f"word_btn_{i}"

                col1, col2 = st.columns([6, 6]) # Adjust column width for button placement

                with col1:
                    if st.button("🔊 Read Aloud (Summarized)", key=button_key_vocalize):
                        if not groq_api_key:
                            st.error("Groq API key is not set. Cannot generate audio.")
                        else:
                            with st.spinner("Summarizing and generating audio..."):
                                try:
                                    # 1. Summarize the response
                                    summarization_service_chat = SummarizationService(api_choice="groq", groq_api_key=groq_api_key)
                                    summarized_text = summarization_service_chat.summarize(message["content"], max_length=max_summary_length_chat)
                                    st.info(f"Summarized for audio: {summarized_text}") # Optional: show the summarized text

                                    # 2. Convert the summary to speech
                                    tts_service_chat = TextToSpeechService(groq_api_key=groq_api_key)
                                    if selected_voice_chat:
                                        tts_service_chat.set_voice(selected_voice_chat)
                                    audio_data_chat = tts_service_chat.generate_speech(summarized_text, enhance_text=enhance_text_for_tts_chat)

                                    # 3. Play the audio
                                    if audio_data_chat:
                                        st.audio(audio_data_chat, format='audio/mp3')
                                    else:
                                        st.error("Failed to generate audio.")
                                except ValueError as e:
                                    st.error(f"Configuration Error for audio generation: {e}")
                                except Exception as e:
                                    st.error(f"An unexpected error occurred during audio generation: {e}")

                with col2:
                     # Add button to export individual response to Word
                    word_file_buffer = create_word_from_text(message["content"])
                    st.download_button(
                        label="📄 Export to Word",
                        data=word_file_buffer,
                        file_name=f"chatbot_response_{i+1}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key=button_key_word # Unique key for the download button
                    )


    # Chat input
    if prompt := st.chat_input("Ask me anything..."):
        # Add user message to chat history
        st.session_state.messages.append({"role": "user", "content": prompt})
        # Display user message
        with st.chat_message("user"):
            st.markdown(prompt)

        # Get and display assistant response
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                response = get_groq_response(prompt, groq_api_key, model=chatbot_model)
                st.markdown(response)
                # Add assistant response to chat history
                st.session_state.messages.append({"role": "assistant", "content": response})

                # Add the vocalize and export buttons after the response is displayed
                # Use the index of the newly added message
                i = len(st.session_state.messages) - 1
                button_key_vocalize = f"vocalize_btn_{i}"
                button_key_word = f"word_btn_{i}"

                col1, col2 = st.columns([6, 6]) # Adjust column width for button placement

                with col1:
                    if st.button("🔊 Read Aloud (Summarized)", key=button_key_vocalize):
                         if not groq_api_key:
                            st.error("Groq API key is not set. Cannot generate audio.")
                         else:
                            with st.spinner("Summarizing and generating audio..."):
                                try:
                                    # 1. Summarize the response
                                    summarization_service_chat = SummarizationService(api_choice="groq", groq_api_key=groq_api_key)
                                    summarized_text = summarization_service_chat.summarize(response, max_length=max_summary_length_chat)
                                    st.info(f"Summarized for audio: {summarized_text}") # Optional: show the summarized text

                                    # 2. Convert the summary to speech
                                    tts_service_chat = TextToSpeechService(groq_api_key=groq_api_key)
                                    if selected_voice_chat:
                                        tts_service_chat.set_voice(selected_voice_chat)
                                    audio_data_chat = tts_service_chat.generate_speech(summarized_text, enhance_text=enhance_text_for_tts_chat)

                                    # 3. Play the audio
                                    if audio_data_chat:
                                        st.audio(audio_data_chat, format='audio/mp3')
                                    else:
                                        st.error("Failed to generate audio.")
                                except ValueError as e:
                                    st.error(f"Configuration Error for audio generation: {e}")
                                except Exception as e:
                                    st.error(f"An unexpected error occurred during audio generation: {e}")
                with col2:
                    # Add button to export individual response to Word
                    word_file_buffer = create_word_from_text(response)
                    st.download_button(
                        label="📄 Export to Word",
                        data=word_file_buffer,
                        file_name=f"chatbot_response_{i+1}.docx",
                        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        key=button_key_word # Unique key for the download button
                    )


st.markdown("---")
st.markdown("Built with Streamlit, Hugging Face, and Groq APIs.")
